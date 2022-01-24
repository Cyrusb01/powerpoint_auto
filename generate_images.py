"""
This script edits the powerpoint and saves it as three seperate images
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.chart.data import CategoryChartData
from data_collection import connect_to_sheet
import os 
from datetime import datetime, timedelta
import pandas as pd
import win32com.client


def replace_paragraph_text_retaining_initial_formatting(paragraph, new_text):
    p = paragraph._p  # the lxml element containing the `<a:p>` paragraph element
    # remove all but the first run
    for idx, run in enumerate(paragraph.runs):
        if idx == 0:
            continue
        p.remove(run._r)
    paragraph.runs[0].text = new_text




prs=Presentation("BFC_Twitter_Post.pptx")
line_chart_df, drawdown_chart_df, performance_table, holdings, weights, strategy, strategy_weights, hedge, hedge_weights = connect_to_sheet()

######################## SLIDE 1 ########################
slide = prs.slides[0]
#Loop through all the elements of Slide 1
for shape in slide.shapes:
    if not shape.has_text_frame:
        
        #Line chart on First page 
        if shape.has_chart:
            line_chart = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = line_chart_df.index.to_list()
            chart_data.add_series('BFC Multi-Strat Fund',  line_chart_df.iloc[:, 0].to_list())
            chart_data.add_series('Bitcoin',    line_chart_df.iloc[:, 1].to_list())
            line_chart.replace_data(chart_data)
            # print("has chart")
        
        #Performance Table
        if shape.has_table:
            for row in range(1, 2):
                for col in range(6):
                    paragraph = shape.table.cell(row,col).text_frame.paragraphs[0]
                    replace_paragraph_text_retaining_initial_formatting(paragraph, performance_table[row][col])
            # print("has table")
        continue
    
    
    
    text_frame = shape.text_frame
    if "Data" == text_frame.text[:4]:
        citing = text_frame.text

        # date_range = citing[citing.index('as of '): citing.index(". Please")]
        today = (datetime.now() - timedelta(days=1)).strftime("%m/%d/%Y")
        new_citing = citing[:citing.index('as of')+6] + today +  citing[citing.index(". Please"):]
        # print(new_citing)
        paragraph = text_frame.paragraphs[0]
        replace_paragraph_text_retaining_initial_formatting(paragraph, new_citing)
        
######################### SLIDE 2 ########################
slide = prs.slides[1]
#Loop through all the elements of Slide 2
for shape in slide.shapes:

    if shape.has_chart:

        #Drawdown Chart
        if ([s.name for s in shape.chart.series][0]) == "Bitcoin":
            drawdown = shape.chart
            chart_data = CategoryChartData()
            chart_data.categories = drawdown_chart_df.index.to_list()
            chart_data.add_series('Bitcoin',    drawdown_chart_df.iloc[:, 1].to_list())
            chart_data.add_series('BFC Multi-Strat Fund',  drawdown_chart_df.iloc[:, 0].to_list())
            
            drawdown.replace_data(chart_data)
            # print("has chart")
        
        #Pie Chart
        if ([s.name for s in shape.chart.series][0]) == "Weights":
            # print(strategy_weights)
            pie = shape.chart
            pie_data = CategoryChartData()
            pie_data.categories = strategy[1:]
            pie_data.add_series('Series 1', strategy_weights)

            pie.replace_data(pie_data)
            # print("Pie")
        #Bar Chart
        if ([s.name for s in shape.chart.series][0]) == "Hedge":
            bar = shape.chart
            bar_data = CategoryChartData()
            bar_data.categories = hedge[1:]
            bar_data.add_series('Series 1', hedge_weights)
            bar.replace_data(bar_data)
            # print("bar")
    if shape.has_table:
        for row in range(1, len(holdings)):
            try:
                paragraph = shape.table.cell(row,0).text_frame.paragraphs[0]
                replace_paragraph_text_retaining_initial_formatting(paragraph, holdings[row])
                
                paragraph = shape.table.cell(row,1).text_frame.paragraphs[0]
                replace_paragraph_text_retaining_initial_formatting(paragraph, weights[row])
            except:
                paragraph = shape.table.cell(0,0).text_frame.paragraphs[0]
                # replace_paragraph_text_retaining_initial_formatting(paragraph, holdings[row])
                extra = True

        
        # print("has table")
    
    if shape.has_text_frame:
        text_frame = shape.text_frame
        # print(text_frame.text)
        if "Data" == text_frame.text[:4]:
            citing = text_frame.text

            # date_range = citing[citing.index('as of '): citing.index(". Please")]
            today = (datetime.now() - timedelta(days=1)).strftime("%m/%d/%Y")
            new_citing = citing[:citing.index('as of')+6] + today +  citing[citing.index(". Please"):citing.index("page.")]
            paragraph = text_frame.paragraphs[0]
            replace_paragraph_text_retaining_initial_formatting(paragraph, new_citing)
    # if not shape.has_text_frame:
        
    #     #Line chart on First page 
    #     if shape.has_chart:
    #         line_chart = shape.chart
    #         chart_data = CategoryChartData()
    #         chart_data.categories = line_chart_df.index.to_list()
    #         chart_data.add_series('BFC Multi-Strat Fund',  line_chart_df.iloc[:, 0].to_list())
    #         chart_data.add_series('Bitcoin',    line_chart_df.iloc[:, 1].to_list())
    #         line_chart.replace_data(chart_data)
    #         print("has chart")
        
    #     #Performance Table
    #     if shape.has_table:
    #         for row in range(1, 2):
    #             for col in range(6):
    #                 paragraph = shape.table.cell(row,col).text_frame.paragraphs[0]
    #                 replace_paragraph_text_retaining_initial_formatting(paragraph, performance_table[row][col])
    #         print("has table")
    #     continue
    
    
    
    # text_frame = shape.text_frame
    # print(text_frame.text[:4])
    # if "Data" == text_frame.text[:4]:
    #     citing = text_frame.text

    #     # date_range = citing[citing.index('as of '): citing.index(". Please")]
    #     today = datetime.now().strftime("%m/%d/%Y")
    #     new_citing = citing[:citing.index('as of')+6] + today +  citing[citing.index(". Please"):]
    #     print(new_citing)
    #     paragraph = text_frame.paragraphs[0]
    #     replace_paragraph_text_retaining_initial_formatting(paragraph, new_citing)
        



prs.save("pptfile.ppt")


Application = win32com.client.Dispatch("PowerPoint.Application")

Presentation = Application.Presentations.Open(os.getcwd() + r"\pptfile.ppt", WithWindow=False)
Presentation.Slides[0].Export(os.getcwd() + r"\slide_images\slide_1.png", "PNG", 2700, 1404)
Presentation.Slides[1].Export(os.getcwd() + r"\slide_images\slide_2.png", "PNG", 2700, 1404)
Presentation.Slides[2].Export(os.getcwd() + r"\slide_images\slide_3.png", "PNG", 2700, 1404)
Application.Quit()
Presentation =  None
Application = None

